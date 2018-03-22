from pptx import Presentation


present=Presentation("test.pptx")
for slide in present.slides:
    if not slide.shapes.title.has_text_frame:
        continue
    print slide.shapes.title.text
    print '................'
    search=GoogleSearch().search(slide.shapes.title.text+" IR")
    for result in search.results:
        print("Title: " + result.title)
        print("Content: " + result.url)
        #if not shape.has_text_frame:
        #    continue
        #title = shape.title
    print '------------------------------------------------------'