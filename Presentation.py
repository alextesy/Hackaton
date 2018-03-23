from __future__ import print_function

import os
from os import listdir
from os.path import isfile, join
from random import randint

from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Cm
from pptx.util import Pt

from Parser import Parser
from Slide import Slide

"""
This class is controlling the powerpoint file. 
It can read the powerpoint text and get the titles of each slide.
Also it will reformat the slides to free space for images/meme and new content.
"""


def smaller_text(slide, iNum):
    """
    Change font size of the slide for all text beside the title
    :param slide: the slide object we want to edit
    """
    for shape in slide.shapes:
        if not shape.has_text_frame or iNum == 0:
            continue

        if shape.is_placeholder:
            if shape.placeholder_format.type is PP_PLACEHOLDER_TYPE.TITLE or \
                    shape.placeholder_format.type is PP_PLACEHOLDER_TYPE.CENTER_TITLE:
                continue
        for paragraph in shape.text_frame.paragraphs:

            if iNum % 3 == 1:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            elif iNum % 3 == 2:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            else:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
            for run in paragraph.runs:
                run.font.size = Pt(15)


def insert_image(slide, path, iNum):
    if iNum == 0:
        pass
    else:
        if iNum % 3 == 2:
            top = Cm(11)
            left = Cm(14)
        elif iNum % 3 == 0:
            top = Cm(11)
            left = Cm(0)
        else:
            top = Cm(11)
            left = Cm(0)

        height = Cm(6.8)
        slide.shapes.add_picture(path, left, top, height=height)


class Pres(object):

    def __init__(self, path, course):
        my_file = Path(path)
        if my_file.is_file() and (path.endswith('.pptx')):
            self.prs = Presentation(path)
            self.path = path
            self.course = course
            self.slides = []
            self.parser = Parser()

        else:
            raise Exception("No Presentation is found")

    def initialize(self):
        for slide in self.prs.slides:
            if slide.shapes.title is None or (not slide.shapes.title.has_text_frame):
                continue
            title = self.parser.parse_title(slide.shapes.title.text)
            s = Slide(title, "", self.prs.slides.index(slide))
            self.slides.append(s)

    def create_url_file(self):
        dir, file = os.path.split(self.path)
        url_file_name = os.path.join(dir, "url_links.txt")
        url_file = open(url_file_name, "w")
        for slide in self.slides:
            print(slide.url.path, file=url_file)

    def create_new_presentation(self):

        for slide in self.prs.slides:
            slide_num = self.prs.slides.index(slide)
            smaller_text(slide, slide_num)

            self.add_image("s" + str(slide_num), slide, slide_num)
        self.add_notes()
        self.add_article()
        self.prs.save("result.pptx")

    def add_notes(self):
        for slide in self.slides:
            slide_note = self.prs.slides[slide.slideNum].notes_slide
            text_frame = slide_note.notes_text_frame
            if text_frame is None:
                text_frame = slide_note.add_text_frame()
            for url in slide.url:
                p = text_frame.add_paragraph()
                r = p.add_run()
                r.text = url
                hlink = r.hyperlink
                hlink.address = url

    def add_image(self, path_to_image_dir, slide, iNum):
        if self.prs.slides.index(slide) % 2 == 0:
            title = self.slides[self.prs.slides.index(slide)].title
            if title.endswith(" "):
                title = title[:-1]
            path_to_image_dir += "\\" + title
        else:
            title = self.slides[self.prs.slides.index(slide)].title
            if title.endswith(" "):
                title = title[:-1]
            path_to_image_dir += "\\" + title + " meme"
        only_files = [f for f in listdir(path_to_image_dir) if isfile(join(path_to_image_dir, f))]
        index = randint(0, len(only_files) - 1)

        path_to_image = path_to_image_dir + "\\" + only_files[index]
        insert_image(slide, path_to_image, iNum)

    def add_article(self):
        blank_slide_layout = self.prs.slide_layouts[6]
        images = [f for f in listdir("Articles") if isfile(join("Articles", f))]
        for image in images:
            img_path = "Articles/" + image
            left = Cm(6)
            top = 0
            height = self.prs.slide_height

            slide = self.prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
